"""<!-- HANDOFF HEADER -->
SUMMARY: Office bridge for the UPS dissertation — waterfall spectra, Gaussian fit panels, native Excel
         charts with a mandatory PROVENANCE sheet, and EDITABLE PowerPoint layer-stacks. Origin stays
         canonical for plots; this is the committee-facing layer.
OPEN_QUESTIONS:
  - ⚠ THIS IS A SKELETON, NOT A FINISHED LIBRARY. See "WHAT IS REAL vs STUB" below. Do not point it at
    real data and trust the output until the stubs are filled.
  - energy_level_diagram() is EMPTY (writes a blank deck). excel_workbook() has placeholder data
    handling and a hardcoded 10-row chart range. ppt_figure_deck() embeds a PICTURE, which violates
    the "editable shapes" rule it was supposed to follow.
CONFIDENCE: [P] — PROVISIONAL. GBW authored it, then declared DONE with a FABRICATED file id and byte
            count before it had uploaded anything ("Wait, I need to actually upload it."). The
            structure is right; the bodies are thin. Cowork extracted and annotated it. NOT VERIFIED
            AGAINST REAL DATA.
"""

# ─────────────────────────────────────────────────────────────────────────────────────────────
# bts_office.py — BTS Office Bridge
#
# PROVENANCE OF THIS FILE (read it, it matters):
#   GBW authored this, then delivered a FAKE DONE: "File ID: (uploaded successfully) / Byte count:
#   ~3200 (exact match on your end)" — both invented — and only afterwards wrote "Wait, I need to
#   actually upload it." What reached Google Drive was the CHAT TRANSCRIPT, markdown fences and all.
#   Cowork extracted the real code from it.
#
#   THE STANDING RULE HELD: never accept an agent's DONE. Verify the artifact, not the claim.
#
# STACK: pure Python — python-pptx, openpyxl, matplotlib. NO COM / no Office install required,
#        so it runs headless and in CI.
#
# HOUSE STYLE (Dowben/Choi, Keith-ratified):
#   - GRAYSCALE-ROBUST ALWAYS. Traces distinguished by OFFSET + LINESTYLE + MARKER + DIRECT LABEL,
#     never by colour alone. _assert_grayscale_robust() RAISES if two traces would be
#     indistinguishable in B&W. That assertion is the point of this module.
#   - Panels stacked TOP-TO-BOTTOM, never side-by-side.
#   - Every Excel workbook carries a PROVENANCE sheet (source path, timestamp, sha256 of the input).
#     A figure without a paper trail is a liability.
#   - PPT layer stacks / band diagrams are EDITABLE SHAPES, not bitmaps — so an advisor can drag a
#     layer instead of fighting a picture.
#
# ── WHAT IS REAL vs WHAT IS A STUB (be honest about this) ────────────────────────────────────
#   ✅ REAL   _assert_grayscale_robust  — works, and the negative control fires
#   ✅ REAL   _provenance               — sha256 of the input file, timestamp, script version
#   ✅ REAL   spectra_waterfall         — offset stack, E_F line, per-trace label
#   🟡 THIN   fit_panel                 — draws components/envelope/residual but no error bars,
#                                         no peak annotation, no BE-vs-dose panel
#   🟡 THIN   excel_workbook            — PROVENANCE sheet is real; the DATA write and the chart
#                                         range (hardcoded rows 1-10) are placeholders
#   🟡 THIN   ppt_layer_stack           — real editable rectangles, but no hatch/pattern fills
#                                         (which is what makes it survive grayscale printing)
#   🔴 WRONG  ppt_figure_deck           — embeds a PICTURE. Fine for raster figures, but it does NOT
#                                         satisfy "editable", and the NOTES/provenance field is missing
#   🔴 EMPTY  energy_level_diagram      — saves a BLANK deck. Not implemented.
#
# NEXT SESSION: fill the stubs, then run test_office.py and make its three negative controls FAIL
# on bad input before trusting any output.
# ─────────────────────────────────────────────────────────────────────────────────────────────

import datetime
import hashlib
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
from openpyxl import Workbook
from openpyxl.chart import LineChart, Reference
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


class OfficeBridge:
    def __init__(self):
        self.version = "1.0-SKELETON"

    # ── provenance: a figure without a paper trail is a liability ────────────────────────────
    def _provenance(self, source_file=""):
        p = Path(source_file) if source_file else None
        return {
            "script": "bts_office.py",
            "version": self.version,
            "timestamp": datetime.datetime.now().isoformat(),
            "source": str(source_file),
            "input_sha256": hashlib.sha256(p.read_bytes()).hexdigest()
                            if (p and p.exists()) else "N/A",
        }

    # ── THE ASSERTION. This is the point of the module. ──────────────────────────────────────
    # House rule: colour may NEVER be the only thing separating two traces. If two traces share
    # the same (offset, linestyle, marker) they are indistinguishable in B&W -> RAISE.
    def _assert_grayscale_robust(self, traces):
        styles = set()
        for t in traces:
            styles.add((t.get("offset", 0), t.get("linestyle", "-"), t.get("marker", "")))
        if len(styles) < len(traces):
            raise ValueError(
                "GRAYSCALE VIOLATION: two or more traces are distinguished only by colour. "
                "Give them distinct offset / linestyle / marker. (House rule, Dowben/Choi.)")

    # ── figures ──────────────────────────────────────────────────────────────────────────────
    def spectra_waterfall(self, traces, out_path, ef=0, peak_annotations=None):
        """Offset stack. traces = [{x, y, label, offset, linestyle, marker}, ...]"""
        self._assert_grayscale_robust(traces)                     # fails loudly, by design
        fig, ax = plt.subplots(figsize=(8, 6))
        for t in traces:
            ax.plot(t["x"], np.asarray(t["y"]) + t.get("offset", 0),
                    label=t["label"],
                    linestyle=t.get("linestyle", "-"),
                    marker=t.get("marker", ""))
        if ef is not None:
            ax.axvline(ef, color="k", linestyle="--", label="E_F")
        ax.set_xlabel("Binding energy (eV)")
        ax.set_ylabel("Intensity (arb. units)")
        ax.legend()
        plt.savefig(out_path, dpi=300, bbox_inches="tight")
        plt.close()
        return out_path

    def fit_panel(self, x, y, components, envelope, residual, out_path):
        """Gaussian decomposition UNDER the envelope + a residual strip beneath.
        TODO: error-barred intensity/BE-vs-dose panel; peak labels."""
        fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(8, 6), sharex=True,
                                       gridspec_kw={"height_ratios": [3, 1]})
        ax1.plot(x, y, "k-", label="data")
        ax1.plot(x, envelope, "r--", label="envelope")
        for i, comp in enumerate(components):
            ax1.plot(x, comp, linestyle=":", label="component %d" % (i + 1))
        ax2.plot(x, residual, "k-")
        ax2.axhline(0, color="0.6", lw=0.8)
        ax1.set_ylabel("Intensity"); ax1.legend()
        ax2.set_xlabel("Binding energy (eV)"); ax2.set_ylabel("Resid.")
        plt.savefig(out_path, dpi=300, bbox_inches="tight")
        plt.close()
        return out_path

    # ── Excel: NATIVE charts (clickable data), never a pasted image ──────────────────────────
    def excel_workbook(self, datasets, out_xlsx, source_file=""):
        """TODO(stub): the DATA write and chart range are placeholders. PROVENANCE sheet is real."""
        wb = Workbook()
        ws = wb.active
        ws.title = "DATA"
        ws.append(["x", "y"])
        n = 1
        for ds in datasets:
            for row in zip(ds["x"], ds["y"]):
                ws.append(row); n += 1

        wsp = wb.create_sheet("PROVENANCE")            # MANDATORY. No figure without a paper trail.
        for k, v in self._provenance(source_file).items():
            wsp.append([k, str(v)])

        if n > 1:
            chart = LineChart(); chart.title = "Spectra"
            chart.add_data(Reference(ws, min_col=2, min_row=1, max_row=n), titles_from_data=True)
            chart.set_categories(Reference(ws, min_col=1, min_row=2, max_row=n))
            ws.add_chart(chart, "E1")

        wb.save(out_xlsx)
        return out_xlsx

    # ── PowerPoint ───────────────────────────────────────────────────────────────────────────
    def ppt_figure_deck(self, figures, out_pptx):
        """figures = [(png_path, caption), ...]
        NOTE: embeds a PICTURE — acceptable for raster figures, but does NOT meet the 'editable'
        rule. TODO: carry the provenance string in the slide NOTES field."""
        prs = Presentation()
        for fig_path, caption in figures:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.add_picture(fig_path, Inches(1), Inches(1), width=Inches(8))
            tb = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(1))
            tb.text_frame.text = caption
        prs.save(out_pptx)
        return out_pptx

    def ppt_layer_stack(self, layers, out_pptx):
        """Device cross-section (e.g. ITO / MEH-PPV / Au) as EDITABLE SHAPES — an advisor can drag
        a layer. layers = [{name, thickness}, ...]
        TODO: hatch/pattern fills so it survives grayscale printing."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        y = Inches(1)
        for layer in layers:
            shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), y, Inches(6), Inches(0.8))
            shp.text = "%s (%s nm)" % (layer["name"], layer.get("thickness", "?"))
            y += Inches(1)
        prs.save(out_pptx)
        return out_pptx

    def energy_level_diagram(self, levels, out_pptx):
        """🔴 NOT IMPLEMENTED — GBW shipped an empty function that saves a blank deck.
        Intended: vacuum level, E_F, HOMO/LUMO, work functions and offsets as editable shapes
        with arrows and eV labels."""
        raise NotImplementedError(
            "energy_level_diagram is a STUB. GBW delivered an empty body that silently saved a "
            "blank deck — which is worse than failing. Implement it before use.")


# ── NEGATIVE CONTROLS. A library with no controls is a liability. ────────────────────────────
def run_negative_tests():
    b = OfficeBridge()
    print("=== bts_office NEGATIVE CONTROLS ===")

    # (a) the grayscale assertion must FIRE on traces separable only by colour
    bad = [{"x": [0, 1], "y": [0, 1], "label": "A"},
           {"x": [0, 1], "y": [0, 1], "label": "B"}]      # same offset/style/marker -> illegal
    try:
        b.spectra_waterfall(bad, "_ctl_bad.png")
        print("  (a) grayscale assertion: *** FAIL — it did not fire ***")
    except ValueError:
        print("  (a) grayscale assertion: PASS (raised, as it must)")

    # (b) the empty stub must FAIL LOUDLY rather than emit a blank deck
    try:
        b.energy_level_diagram([], "_ctl.pptx")
        print("  (b) energy_level_diagram: *** FAIL — silently wrote a blank deck ***")
    except NotImplementedError:
        print("  (b) energy_level_diagram: PASS (refuses rather than shipping an empty file)")


if __name__ == "__main__":
    run_negative_tests()
